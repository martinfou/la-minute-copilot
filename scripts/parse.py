#!/usr/bin/env python3
"""
Parser PowerPoint pour La minute Copilot.
Lit un fichier JSON de contenu et génère un .pptx.

Format du JSON:
  slides[].title          → Titre de la diapositive
  slides[].subtitle       → Sous-titre
  slides[].bullets[]      → Points clés (affichés sur la slide)
  slides[].speaker_notes  → Notes du présentateur (cachées dans le PPT)
  slides[].cards[]        → Cartes visuelles (optionnelles)
  slides[].footer         → Pied de page
  slides[].bonus          → Bannière bonus (summary)

Usage:
    python3 parse.py presentations/ep01-outlook.json
    python3 parse.py presentations/ep01-outlook.json -o mon_fichier.pptx
"""

import json, sys, os, argparse
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Couleurs ──
BLUE = RGBColor(0x00, 0x78, 0xD4)
DARK = RGBColor(0x2B, 0x2B, 0x2B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x6B, 0x6B, 0x6B)
LIGHT_BG = RGBColor(0xF0, 0xF4, 0xF8)
GREEN = RGBColor(0x10, 0xB9, 0x81)
AMBER = RGBColor(0xF5, 0x9E, 0x0B)
CARD_BORDER = RGBColor(0xE0, 0xE0, 0xE0)
FOOTER_BLUE = RGBColor(0xAA, 0xD4, 0xF0)
BONUS_BG = RGBColor(0xFF, 0xF3, 0xCD)
BONUS_TEXT = RGBColor(0x85, 0x6E, 0x04)

PALETTE = {'blue': BLUE, 'green': GREEN, 'amber': AMBER}

def add_bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color

def add_tb(slide, l, t, w, h, text, size=18, bold=False, color=DARK, align=PP_ALIGN.LEFT):
    """Ajoute une textbox."""
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = 'Calibri'
    p.alignment = align
    return tf

def add_lines(slide, l, t, w, h, lines, size=13, color=DARK, bullet=False):
    """Ajoute des lignes de texte."""
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = 'Calibri'
        p.space_after = Pt(4)
        if bullet and line.strip() and not line.strip().startswith('  '):
            p.level = 0
            p.font.bold = True
        elif bullet and line.startswith('  '):
            p.level = 1
            p.font.size = Pt(size - 1)
    return tf

def add_card(slide, l, t, w, h, card):
    """Carte avec fond blanc, bordure, contenu."""
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = CARD_BORDER

    accent = PALETTE.get(card.get('color', 'blue'), BLUE)
    add_tb(slide, l + 0.3, t + 0.2, w - 0.6, 0.5,
           f"{card.get('icon', '')} {card.get('title', '')}", size=18, bold=True, color=accent)
    add_lines(slide, l + 0.3, t + 0.7, w - 0.6, h - 1.0, card.get('lines', []), size=13, color=DARK)

def generate(json_path, output_path=None):
    with open(json_path, 'r', encoding='utf-8') as f:
        data = json.load(f)

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    for sd in data['slides']:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        stype = sd.get('type', '')

        # Speaker notes
        if sd.get('speaker_notes'):
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = sd['speaker_notes']

        if stype == 'title':
            add_bg(slide, BLUE)
            add_tb(slide, 1, 2.0, 11, 1.2, sd['title'], size=48, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
            add_tb(slide, 1, 3.2, 11, 0.8, sd.get('subtitle', ''), size=28, color=FOOTER_BLUE, align=PP_ALIGN.CENTER)
            add_tb(slide, 1, 4.5, 11, 0.6, sd.get('footer', ''), size=16, color=FOOTER_BLUE, align=PP_ALIGN.CENTER)

        elif stype == 'tips':
            add_bg(slide, WHITE)
            add_tb(slide, 0.5, 0.3, 12, 0.7, sd['title'], size=32, bold=True, color=DARK)
            add_tb(slide, 0.5, 1.0, 12, 0.5, sd.get('subtitle', ''), size=18, color=GRAY)

            # Bullet points
            bullets = sd.get('bullets', [])
            if bullets:
                x_bullets = 0.5
                w_bullets = 12.0
            
            # Cards
            cards = sd.get('cards', [])
            if cards:
                n = len(cards)
                cw = (12.0 / n) - 0.3
                x = 0.5
                top = 1.7
                
                # If bullets exist, put them above cards
                if bullets:
                    add_lines(slide, 0.5, 1.7, 12.0, 1.5, [f"• {b}" for b in bullets], size=14, color=DARK)
                    top = 3.3
                
                for card in cards:
                    add_card(slide, x, top, cw, 4.0, card)
                    x += cw + 0.3
            elif bullets:
                add_lines(slide, 0.5, 1.7, 12.0, 5.0, [f"• {b}" for b in bullets], size=16, color=DARK)

        elif stype == 'bonus':
            add_bg(slide, LIGHT_BG)
            add_tb(slide, 0.5, 0.3, 12, 0.7, sd['title'], size=32, bold=True, color=DARK)

            bullets = sd.get('bullets', [])
            cards = sd.get('cards', [])

            if bullets and not cards:
                add_lines(slide, 0.5, 1.3, 12.0, 5.5, [f"• {b}" for b in bullets], size=16, color=DARK)
            elif cards:
                n = len(cards)
                cw = (12.0 / n) - 0.3
                x = 0.5
                for card in cards:
                    add_card(slide, x, 1.3, cw, 5.5, card)
                    x += cw + 0.3

        elif stype == 'summary':
            add_bg(slide, BLUE)
            add_tb(slide, 1, 1.0, 11, 0.8, sd['title'], size=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

            cards = sd.get('cards', [])
            n = len(cards)
            if n > 0:
                cw = 3.8
                gap = (11 - n * cw) / (n + 1)
                x = 0.5 + gap
                for card in cards:
                    shape = slide.shapes.add_shape(1, Inches(x), Inches(2.2), Inches(cw), Inches(2.5))
                    shape.fill.solid()
                    shape.fill.fore_color.rgb = WHITE
                    shape.line.color.rgb = FOOTER_BLUE
                    add_tb(slide, x + 0.3, 2.4, cw - 0.6, 0.6,
                            f"{card.get('icon', '')} {card.get('title', '')}", size=22, bold=True, color=DARK, align=PP_ALIGN.CENTER)
                    add_tb(slide, x + 0.3, 3.2, cw - 0.6, 1.0,
                            card.get('desc', ''), size=16, color=GRAY, align=PP_ALIGN.CENTER)
                    x += cw + gap

            bonus = sd.get('bonus', '')
            if bonus:
                shape = slide.shapes.add_shape(1, Inches(3.5), Inches(5.2), Inches(6.3), Inches(1.2))
                shape.fill.solid()
                shape.fill.fore_color.rgb = BONUS_BG
                shape.line.color.rgb = RGBColor(0xFF, 0xE0, 0x82)
                add_tb(slide, 3.7, 5.4, 5.9, 0.6, bonus, size=16, bold=True, color=BONUS_TEXT, align=PP_ALIGN.CENTER)

            footer = sd.get('footer', '')
            if footer:
                add_tb(slide, 1, 6.7, 11, 0.5, footer, size=14, color=FOOTER_BLUE, align=PP_ALIGN.CENTER)

    # Save
    if not output_path:
        basename = os.path.splitext(os.path.basename(json_path))[0]
        output_path = os.path.join(os.path.dirname(json_path), f"{basename}.pptx")

    prs.save(output_path)
    print(f"✅ Présentation générée: {output_path}")
    print(f"   Diapositives: {len(prs.slides)}")
    print(f"   Taille: {os.path.getsize(output_path) / 1024:.1f} KB")

if __name__ == '__main__':
    p = argparse.ArgumentParser(description='Génère un PowerPoint depuis un JSON')
    p.add_argument('json', help='Fichier JSON des slides')
    p.add_argument('-o', '--output', help='Fichier .pptx de sortie')
    args = p.parse_args()
    generate(args.json, args.output)
