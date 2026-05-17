from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

BLUE = RGBColor(0x00, 0x78, 0xD4)
DARK = RGBColor(0x2B, 0x2B, 0x2B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x6B, 0x6B, 0x6B)
LIGHT_BG = RGBColor(0xF5, 0xF5, 0xF5)
GREEN = RGBColor(0x10, 0xB9, 0x81)

def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_textbox(slide, left, top, width, height, text, size=18, bold=False, color=DARK, align=PP_ALIGN.LEFT, font_name='Calibri'):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = align
    return tf

def add_multiline(slide, left, top, width, height, lines, size=16, color=DARK, spacing=1.5, font_name='Calibri'):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = font_name
        p.space_after = Pt(6)
    return tf

def add_card(slide, left, top, width, height, title, content, icon="💡", accent=BLUE):
    from pptx.util import Emu
    shape = slide.shapes.add_shape(
        1, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.shadow.inherit = False
    shape.line.color.rgb = RGBColor(0xE0, 0xE0, 0xE0)
    
    # Icon + title
    add_textbox(slide, left + 0.3, top + 0.2, width - 0.6, 0.5, f"{icon} {title}", size=18, bold=True, color=accent)
    
    # Content
    lines = content.split('\n')
    add_multiline(slide, left + 0.3, top + 0.7, width - 0.6, height - 1.0, lines, size=13, color=DARK)

# ── Slide 1: Title ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BLUE)
add_textbox(slide, 1, 2.0, 11, 1.2, "🤖 La minute Copilot", size=48, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_textbox(slide, 1, 3.2, 11, 0.8, "Copilot dans Outlook", size=28, bold=False, color=RGBColor(0xCC,0xE5,0xFF), align=PP_ALIGN.CENTER)
add_textbox(slide, 1, 4.5, 11, 0.6, "3 astuces pour IT et non-IT  •  Sauvegarder ses prompts", size=16, color=RGBColor(0xAA,0xD4,0xF0), align=PP_ALIGN.CENTER)

# ── Slide 2: Astuce 1 ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_textbox(slide, 0.5, 0.3, 12, 0.7, "🧠 Astuce 1 — Résumer un fil de discussion", size=32, bold=True, color=BLUE)
add_textbox(slide, 0.5, 1.0, 12, 0.5, "Fini les longues lectures de 15 courriels", size=18, color=GRAY)

add_card(slide, 0.5, 1.7, 4.0, 5.0,
    "Comment faire?",
    "1. Ouvrir un courriel dans un fil\n2. Cliquer sur l'icône Copilot\ndans le ruban\n3. Choisir «Summarize»\n\nOu taper directement votre\ndemande en français!",
    "📋", BLUE)

add_card(slide, 4.7, 1.7, 4.0, 5.0,
    "Exemples de prompts",
    "«Résume ce fil de discussion\net liste les décisions prises»\n\n«Quelles sont les dates de\nlivraison mentionnées?»\n\n«Extrais les action items\navec responsables et\néchéances»",
    "💬", GREEN)

add_card(slide, 8.9, 1.7, 4.0, 5.0,
    "Qui? Pourquoi?",
    "👨‍💼 Gestionnaires:\nSuivre 5 projets à la fois\n\n👩‍💻 IT:\nRetour de congé → résumé\ndes 50 courriels en 30 sec\n\n👤 Non-IT:\nFils interminables de clients",
    "🎯", RGBColor(0xF5,0x9E,0x0B))

# ── Slide 3: Astuce 2 ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_textbox(slide, 0.5, 0.3, 12, 0.7, "✍️ Astuce 2 — Rédiger une réponse en 2 clics", size=32, bold=True, color=BLUE)
add_textbox(slide, 0.5, 1.0, 12, 0.5, "Plus besoin de chercher ses mots — décrivez, Copilot écrit", size=18, color=GRAY)

add_card(slide, 0.5, 1.7, 6.0, 5.0,
    "Comment faire?",
    "1. Cliquez sur «Reply» normalement\n2. Dans le bas de la fenêtre, cliquez sur l'icône Copilot\n\n3. Tapez ce que vous voulez dire:\n«Remercie Jean et confirme que le\n   rapport sera prêt vendredi»\n\n4. Copilot génère un brouillon\n→ Vous éditez ou envoyez",
    "📋", BLUE)

add_card(slide, 6.7, 1.7, 6.0, 5.0,
    "💡 Ajustement de ton (très pratique!)",
    "Après la génération, vous pouvez\nlui demander:\n\n• «Make this more professional»\n• «Make it friendlier»\n• «Make it more concise»\n• «Réécris en ton neutre»\n• «Version courrier interne»\n\n🎯 Pour qui? Tout le monde!\nSurtout ceux qui écrivent\nbeaucoup de courriels par jour",
    "🎨", GREEN)

# ── Slide 4: Astuce 3 ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_textbox(slide, 0.5, 0.3, 12, 0.7, "🎯 Astuce 3 — Prioriser sa boîte de réception", size=32, bold=True, color=BLUE)
add_textbox(slide, 0.5, 1.0, 12, 0.5, "Au lieu de fouiller 50+ courriels, laissez Copilot trier pour vous", size=18, color=GRAY)

add_card(slide, 0.5, 1.7, 4.0, 5.0,
    "Comment faire?",
    "Ouvrez le volet Copilot\n(icône dans le ruban à droite)\n\nPuis écrivez:\n\n«Show me my most urgent\nunread emails from today»\n\nC'est tout!",
    "📋", BLUE)

add_card(slide, 4.7, 1.7, 4.0, 5.0,
    "Variations utiles",
    "«Prioritize my emails this\nweek by importance»\n\n«Quels courriels nécessitent\nune action de ma part\ndemain?»\n\n«Résume les courriels de\nmon gestionnaire des\nderniers jours»",
    "💬", GREEN)

add_card(slide, 8.9, 1.7, 4.0, 5.0,
    "Scénarios gagnants",
    "👨‍💼 Lundi matin: prioriser\n50 courriels du week-end\n\n👩‍💻 Retour de vacances:\ntout savoir en 2 minutes\n\n👤 Fin de mois:\nne rien oublier avant\nla date limite",
    "🎯", RGBColor(0xF5,0x9E,0x0B))

# ── Slide 5: Sauvegarder ses prompts ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, RGBColor(0xF0,0xF4,0xF8))
add_textbox(slide, 0.5, 0.3, 12, 0.7, "💾 Bonus — Sauvegarder ses prompts Copilot", size=32, bold=True, color=DARK)

# Quick Parts card
shape = slide.shapes.add_shape(1, Inches(0.5), Inches(1.3), Inches(5.8), Inches(5.5))
shape.fill.solid()
shape.fill.fore_color.rgb = WHITE
shape.line.color.rgb = RGBColor(0xE0,0xE0,0xE0)
add_textbox(slide, 0.8, 1.5, 5.2, 0.5, "⭐ Quick Parts (recommandé)", size=20, bold=True, color=BLUE)
add_multiline(slide, 0.8, 2.1, 5.2, 4.5, [
    "1. Nouveau courriel → tapez votre prompt",
    "2. Sélectionnez le texte",
    "3. Insertion → Quick Parts →",
    "   Save Selection to Quick Part Gallery",
    "4. Nommez: «[CP] Résumer fil»",
    "",
    "✅ Pour réutiliser:",
    "Insertion → Quick Parts → votre prompt",
    "",
    "💡 Astuce: Créez-en 3-4 pour vos",
    "tâches les plus fréquentes!",
], size=14, color=DARK)

# Other options card
shape2 = slide.shapes.add_shape(1, Inches(6.8), Inches(1.3), Inches(6.0), Inches(2.5))
shape2.fill.solid()
shape2.fill.fore_color.rgb = WHITE
shape2.line.color.rgb = RGBColor(0xE0,0xE0,0xE0)
add_textbox(slide, 7.1, 1.5, 5.5, 0.5, "📄 OneNote / Bloc-Notes", size=18, bold=True, color=GREEN)
add_multiline(slide, 7.1, 2.1, 5.5, 1.5, [
    "Page partagée «Mes prompts Copilot»",
    "Classez par app: Outlook, Teams, Word...",
    "Copier-coller → prompt utilisé en 2 sec",
], size=14, color=DARK)

# Copilot Lab card
shape3 = slide.shapes.add_shape(1, Inches(6.8), Inches(4.2), Inches(6.0), Inches(2.5))
shape3.fill.solid()
shape3.fill.fore_color.rgb = WHITE
shape3.line.color.rgb = RGBColor(0xE0,0xE0,0xE0)
add_textbox(slide, 7.1, 4.4, 5.5, 0.5, "🏪 Copilot Lab (officiel)", size=18, bold=True, color=RGBColor(0xF5,0x9E,0x0B))
add_multiline(slide, 7.1, 5.0, 5.5, 1.5, [
    "Galerie gratuite de prompts Microsoft",
    "copilot.cloud.microsoft/prompts",
    "Centaines de prompts prêts à copier",
], size=14, color=DARK)

# ── Slide 6: Summary ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BLUE)
add_textbox(slide, 1, 1.0, 11, 0.8, "🎯 Les 3 astuces + le bonus", size=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

cards_data = [
    ("🧠", "Résumer", "Fini les longs fils\nde discussion"),
    ("✍️", "Rédiger", "Décrire en 5 mots\nCopilot écrit le reste"),
    ("🎯", "Prioriser", "50 courriels triés\nen 30 secondes"),
]

for i, (icon, title, desc) in enumerate(cards_data):
    x = 0.5 + i * 4.3
    shape = slide.shapes.add_shape(1, Inches(x), Inches(2.2), Inches(3.8), Inches(2.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = RGBColor(0xCC,0xE5,0xFF)
    add_textbox(slide, x + 0.3, 2.4, 3.2, 0.6, f"{icon} {title}", size=22, bold=True, color=DARK, align=PP_ALIGN.CENTER)
    add_textbox(slide, x + 0.3, 3.2, 3.2, 1.0, desc, size=16, color=GRAY, align=PP_ALIGN.CENTER)

# Bonus at bottom
shape = slide.shapes.add_shape(1, Inches(3.5), Inches(5.2), Inches(6.3), Inches(1.2))
shape.fill.solid()
shape.fill.fore_color.rgb = RGBColor(0xFF,0xF3,0xCD)
shape.line.color.rgb = RGBColor(0xFF,0xE0,0x82)
add_textbox(slide, 3.7, 5.4, 5.9, 0.6, "💾 Bonus: Sauvegardez vos prompts avec Quick Parts + OneNote + Copilot Lab", size=16, bold=True, color=RGBColor(0x85,0x6E,0x04), align=PP_ALIGN.CENTER)

add_textbox(slide, 1, 6.7, 11, 0.5, "La minute Copilot — Rendez-vous la semaine prochaine! 👋", size=14, color=RGBColor(0xAA,0xD4,0xF0), align=PP_ALIGN.CENTER)

# Save
output_path = "/home/martinfou/projects/personal-finance/data/La_Minute_Copilot_Outlook.pptx"
os.makedirs(os.path.dirname(output_path), exist_ok=True)
prs.save(output_path)
print(f"✅ PowerPoint sauvegardé: {output_path}")
print(f"   Taille: {os.path.getsize(output_path) / 1024:.1f} KB")
print(f"   Diapositives: {len(prs.slides)}")
