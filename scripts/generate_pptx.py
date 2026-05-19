#!/usr/bin/env python3
"""
Générateur PowerPoint pour La minute Copilot.
Lit un fichier JSON de contenu (avec schéma validé) et génère un .pptx.

Style: thème Desjardins — vert corporatif, fond menthe.
Icônes: Lucide (https://lucide.dev) — glyphes vectoriels monochromes.
Le format JSON est documenté dans schema/la-minute-copilot.schema.json.

Usage:
    python3 scripts/generate_pptx.py presentations/outlook.json
    python3 scripts/generate_pptx.py presentations/outlook.json -o mon_fichier.pptx

Dépendances: python-pptx (pip install python-pptx)
"""

import json, os, argparse, re, sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

sys.path.insert(0, str(Path(__file__).resolve().parent))
from lucide_icons import render_icon_png

# ── Palette Desjardins ──
# Vert corporatif principal (logo Desjardins) — approximation du Pantone 348
DESJARDINS_GREEN = RGBColor(0x00, 0x87, 0x4E)
# Vert plus foncé pour titres sur fond clair
DARK_GREEN = RGBColor(0x00, 0x6B, 0x3F)
# Vert menthe pâle — fond des slides de contenu
LIGHT_GREEN_BG = RGBColor(0xE6, 0xF0, 0xEA)
# Vert pâle pour bordures / accents sur fond vert
PALE_GREEN = RGBColor(0xCC, 0xE0, 0xD2)

# Neutres
DARK = RGBColor(0x2B, 0x2B, 0x2B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x6B, 0x6B, 0x6B)
LIGHT_GRAY = RGBColor(0xE0, 0xE0, 0xE0)

# Couleurs d'accent secondaires
AMBER = RGBColor(0xF5, 0x9E, 0x0B)
RED = RGBColor(0xEF, 0x44, 0x44)

# Palette pour les cartes (clé `color` dans le JSON).
# Le bleu hérité est remappé vers le vert Desjardins pour cohérence visuelle.
PALETTE = {
    'green': DESJARDINS_GREEN,
    'amber': AMBER,
    'blue': DESJARDINS_GREEN,
    'red': RED,
}

FONT = 'Calibri'
ICON_SIZE_CARD = 0.32       # pouces — cartes de contenu
ICON_SIZE_SUMMARY = 0.38    # pouces — cartes récap
_ICON_WARNED = False

# Émojis / pictogrammes — retirés à l'affichage pour un rendu plus professionnel
_EMOJI_RE = re.compile(
    "["
    "\U0001F1E0-\U0001F1FF"
    "\U0001F300-\U0001F5FF"
    "\U0001F600-\U0001F64F"
    "\U0001F680-\U0001F6FF"
    "\U0001F700-\U0001F77F"
    "\U0001F780-\U0001F7FF"
    "\U0001F800-\U0001F8FF"
    "\U0001F900-\U0001F9FF"
    "\U0001FA00-\U0001FAFF"
    "\U00002702-\U000027B0"
    "\U000024C2-\U0001F251"
    "]+",
    flags=re.UNICODE,
)


def clean_text(text):
    """Retire les émojis et normalise les espaces."""
    if not text:
        return text
    text = _EMOJI_RE.sub("", text)
    text = re.sub(r"[\u200d\uFE0F]", "", text)
    return re.sub(r"  +", " ", text).strip()


# ── Primitives ──

def add_bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def add_rect(slide, l, t, w, h, fill_color, line_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line_color
    shape.shadow.inherit = False
    return shape


def add_tb(slide, l, t, w, h, text, size=18, bold=False, color=DARK,
           align=PP_ALIGN.LEFT, font=FONT):
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    p = tf.paragraphs[0]
    p.text = clean_text(text)
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font
    p.alignment = align
    return tf


def add_lines(slide, l, t, w, h, lines, size=13, color=DARK, font=FONT):
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = clean_text(line)
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = font
        p.space_after = Pt(4)
    return tf


def add_lucide_icon(slide, name, l, t, size_inches, color):
    """Insère une icône Lucide rasterisée; retourne False si indisponible."""
    global _ICON_WARNED
    if not name:
        return False
    try:
        size_px = max(24, int(size_inches * 96))
        png_path = render_icon_png(name, color, size_px=size_px)
        slide.shapes.add_picture(str(png_path), Inches(l), Inches(t), width=Inches(size_inches))
        return True
    except Exception as exc:
        if not _ICON_WARNED:
            print(f"⚠️  Icônes Lucide désactivées: {exc}", file=sys.stderr)
            _ICON_WARNED = True
        return False


def add_footer(slide, page_num=None, brand_color=DESJARDINS_GREEN, on_dark=False):
    """Pied de page : trait de séparation, étiquette série, numéro de page."""
    line_color = PALE_GREEN if on_dark else brand_color
    text_color = WHITE if on_dark else GRAY

    add_rect(slide, 0.5, 7.05, 12.33, 0.015, line_color)
    add_tb(slide, 0.5, 7.15, 6, 0.3, "La minute Copilot",
           size=10, color=text_color)

    if page_num is not None:
        add_tb(slide, 12.0, 7.15, 0.9, 0.3, str(page_num),
               size=10, color=text_color, align=PP_ALIGN.RIGHT)


def add_card(slide, l, t, w, h, card):
    """Carte blanche style Desjardins (bordure subtile, icône Lucide + titre accent)."""
    add_rect(slide, l, t, w, h, WHITE, line_color=LIGHT_GRAY)

    accent = PALETTE.get(card.get('color', 'green'), DESJARDINS_GREEN)
    title_l = l + 0.3
    icon_name = card.get('icon', '')
    if icon_name and add_lucide_icon(slide, icon_name, l + 0.28, t + 0.18, ICON_SIZE_CARD, accent):
        title_l = l + 0.28 + ICON_SIZE_CARD + 0.12
    add_tb(slide, title_l, t + 0.2, w - (title_l - l) - 0.2, 0.5,
           card.get('title', ''),
           size=17, bold=True, color=accent)
    add_lines(slide, l + 0.3, t + 0.85, w - 0.6, h - 1.05,
              card.get('lines', []), size=13, color=DARK)


# ── Mises en page ──

def render_title(slide, sd, meta):
    """
    Slide de couverture style Desjardins :
    - Panneau vert plein à gauche (~55%)
    - Titre/sous-titre en blanc, trait de séparation
    - Bas du panneau : auteur, date, footer
    """
    add_bg(slide, LIGHT_GREEN_BG)

    panel_w = 7.4
    add_rect(slide, 0, 0, panel_w, 7.5, DESJARDINS_GREEN)

    # Titre
    add_tb(slide, 0.7, 2.0, panel_w - 1.2, 1.6, sd['title'],
           size=44, bold=True, color=WHITE)

    # Trait de séparation horizontal
    add_rect(slide, 0.7, 3.55, 3.2, 0.04, WHITE)

    # Sous-titre
    add_tb(slide, 0.7, 3.75, panel_w - 1.2, 0.7,
           sd.get('subtitle', ''), size=22, color=WHITE)

    # Bloc info : auteur / date / footer
    info_top = 5.0
    footer = sd.get('footer', '')
    if footer:
        add_tb(slide, 0.7, info_top, panel_w - 1.2, 0.4, footer,
               size=13, color=PALE_GREEN)
        info_top += 0.5

    author = meta.get('author', '')
    if author:
        add_tb(slide, 0.7, info_top, panel_w - 1.2, 0.35,
               author, size=13, bold=True, color=WHITE)
        info_top += 0.35

    date = meta.get('date', '')
    if date:
        add_tb(slide, 0.7, info_top, panel_w - 1.2, 0.3,
               date, size=12, color=PALE_GREEN)

def render_tips(slide, sd, idx):
    """
    Slide de contenu : fond menthe, titre vert,
    puces + cartes blanches, pied de page Desjardins.
    """
    add_bg(slide, LIGHT_GREEN_BG)

    add_tb(slide, 0.5, 0.4, 12, 0.8, sd['title'],
           size=28, bold=True, color=DARK_GREEN)

    subtitle = sd.get('subtitle', '')
    if subtitle:
        add_tb(slide, 0.5, 1.1, 12, 0.45, subtitle, size=15, color=GRAY)

    # Trait de séparation sous le titre
    add_rect(slide, 0.5, 1.65, 12.33, 0.02, DESJARDINS_GREEN)

    bullets = sd.get('bullets', [])
    cards = sd.get('cards', [])

    if cards:
        n = len(cards)
        cw = (12.0 / n) - 0.3
        x = 0.5
        top = 1.9
        card_h = 4.8
        if bullets:
            add_lines(slide, 0.5, 1.9, 12.0, 1.3,
                      [f"• {b}" for b in bullets], size=14, color=DARK)
            top = 3.2
            card_h = 3.6
        for card in cards:
            add_card(slide, x, top, cw, card_h, card)
            x += cw + 0.3
    elif bullets:
        add_lines(slide, 0.5, 1.9, 12.0, 5.0,
                  [f"• {b}" for b in bullets], size=18, color=DARK)

    add_footer(slide, page_num=idx)


def render_bonus(slide, sd, idx):
    """
    Slide bonus : fond menthe, titre vert foncé,
    cartes ou puces, pied de page Desjardins.
    """
    add_bg(slide, LIGHT_GREEN_BG)

    add_tb(slide, 0.5, 0.4, 12, 0.8, sd['title'],
           size=28, bold=True, color=DARK_GREEN)

    add_rect(slide, 0.5, 1.35, 12.33, 0.02, DESJARDINS_GREEN)

    bullets = sd.get('bullets', [])
    cards = sd.get('cards', [])

    if bullets and not cards:
        add_lines(slide, 0.5, 1.6, 12.0, 5.0,
                  [f"• {b}" for b in bullets], size=16, color=DARK)
    elif cards:
        n = len(cards)
        cw = (12.0 / n) - 0.3
        x = 0.5
        for card in cards:
            add_card(slide, x, 1.6, cw, 5.2, card)
            x += cw + 0.3

    add_footer(slide, page_num=idx)


def render_summary(slide, sd, idx):
    """
    Slide finale : fond vert Desjardins plein, titre blanc,
    cartes récap blanches, encadré bonus blanc, footer + logo.
    """
    add_bg(slide, DESJARDINS_GREEN)

    add_tb(slide, 1, 0.8, 11.33, 0.9, sd['title'],
           size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Trait blanc sous le titre
    add_rect(slide, 5.0, 1.85, 3.33, 0.04, WHITE)

    cards = sd.get('cards', [])
    n = len(cards)
    if n > 0:
        available_w = 12.33
        cw = min(3.6, (available_w - 0.3 * (n + 1)) / n)
        gap = (available_w - n * cw) / (n + 1)
        icon_size = ICON_SIZE_SUMMARY if n <= 3 else 0.30
        title_size = 22 if n <= 3 else 17
        desc_size = 15 if n <= 3 else 12
        x = 0.5 + gap
        for card in cards:
            add_rect(slide, x, 2.3, cw, 2.5, WHITE, line_color=PALE_GREEN)
            icon_name = card.get('icon', '')
            title_y = 2.5
            if icon_name:
                icon_x = x + (cw - icon_size) / 2
                if add_lucide_icon(slide, icon_name, icon_x, 2.38, icon_size, DARK_GREEN):
                    title_y = 2.38 + icon_size + 0.08
            add_tb(slide, x + 0.15, title_y, cw - 0.3, 0.6,
                   card.get('title', ''),
                   size=title_size, bold=True, color=DARK_GREEN, align=PP_ALIGN.CENTER)
            desc_y = title_y + 0.7
            add_tb(slide, x + 0.15, desc_y, cw - 0.3, 1.0,
                   card.get('desc', ''), size=desc_size, color=GRAY, align=PP_ALIGN.CENTER)
            x += cw + gap

    bonus = sd.get('bonus', '')
    if bonus:
        # Encadré bonus — même famille visuelle que les cartes récap (blanc + accent vert)
        bx, by, bw, bh = 2.65, 5.12, 8.03, 1.18
        add_rect(slide, bx, by, bw, bh, WHITE, line_color=PALE_GREEN)
        add_rect(slide, bx, by, 0.07, bh, DESJARDINS_GREEN)
        add_tb(slide, bx + 0.18, by + 0.28, bw - 0.36, bh - 0.4, bonus,
               size=15, bold=True, color=DARK_GREEN, align=PP_ALIGN.CENTER)

    footer = sd.get('footer', '')
    if footer:
        add_tb(slide, 1, 6.55, 11.33, 0.4, footer,
               size=13, color=PALE_GREEN, align=PP_ALIGN.CENTER)


# ── Orchestration ──

RENDERERS = {
    'title': render_title,
    'tips': render_tips,
    'bonus': render_bonus,
    'summary': render_summary,
}


def generate(json_path, output_path=None):
    with open(json_path, 'r', encoding='utf-8') as f:
        data = json.load(f)

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    meta = {k: data.get(k, '') for k in ('author', 'date', 'title', 'subtitle')}

    for idx, sd in enumerate(data['slides'], start=1):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        stype = sd.get('type', '')

        if sd.get('speaker_notes'):
            slide.notes_slide.notes_text_frame.text = clean_text(sd['speaker_notes'])

        renderer = RENDERERS.get(stype)
        if renderer is None:
            continue
        if stype == 'title':
            renderer(slide, sd, meta)
        else:
            renderer(slide, sd, idx)

    if not output_path:
        basename = os.path.splitext(os.path.basename(json_path))[0]
        output_path = os.path.join(os.path.dirname(json_path), f"{basename}.pptx")

    prs.save(output_path)
    print(f"✅ Présentation générée: {output_path}")
    print(f"   Diapositives: {len(prs.slides)}")
    print(f"   Taille: {os.path.getsize(output_path) / 1024:.1f} KB")


if __name__ == '__main__':
    p = argparse.ArgumentParser(description='Génère un PowerPoint depuis un JSON (thème Desjardins)')
    p.add_argument('json', help='Fichier JSON des slides')
    p.add_argument('-o', '--output', help='Fichier .pptx de sortie')
    args = p.parse_args()
    generate(args.json, args.output)
